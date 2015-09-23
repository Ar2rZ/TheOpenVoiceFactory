#!/usr/bin/python
"Extracting Utterances from CommuniKate pagesets designed in PowerPoint"
# Make the images export more effectively

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import io
import os
import string
from PIL import Image

alpha = string.ascii_lowercase + string.digits + '_'


def resizeImage(image, scaleFactor):

        oldSize = image.size
        newSize = (scaleFactor*oldSize[0],
                   scaleFactor*oldSize[1])
        return image.resize(newSize, Image.ANTIALIAS)


def remove_punctuation(s):
        """removes puncuation,
        provided by http://codereview.stackexchange.com/a/101806/4759"""
        return ''.join(c for c in s if c in alpha)


def make_title(label):
        """Given a  string, returns the string in the format we use for
         identifying grids. This is used mostly to build internal link
         structures"""
        tag = remove_punctuation(label.lower().strip().replace(" ", "_"))
        if tag == "":
                tag = "unknown"
        return tag


class Locator:

        """Static class designed to abstract away the process of working out
        which bit of the grid a particular part of powerpoint is in"""

        COL_TABLE = {152400: 0, 1503659: 1, 1600200: 1, 2861846: 2,
            2819400: 2, 2854919: 2, 2854925: 2, 4170660: 3,
            4191000: 3, 5542260: 4, 5769114: 4, 5562600: 4, 5769125: 4}
        ROW_TABLE = {0: 0, 152400: 0, 152401: 0, 1981200: 1, 3771900: 2, 5562600: 3,
             5610125: 3, 6095999: 3, 7314625: 4, 7340121: 4, 7340600: 4}

        @staticmethod
        def get_closest_key(dict, inKey):
                "Returns the closest key in the dictionary, for numerical keys."
                # from http://stackoverflow.com/a/7934624/170243
                if inKey in dict:
                        return inKey
                else:
                        return min(dict.keys(), key=lambda k: abs(k - inKey))

        @staticmethod
        def get_cr(topPos, leftPos):
                col = Locator.get_closest_key(Locator.COL_TABLE, leftPos)
                row = Locator.get_closest_key(Locator.ROW_TABLE, topPos)
                return (Locator.COL_TABLE[col], Locator.ROW_TABLE[row])


class Grid:

        """Class representing on n by n grid, complete with utterances, links
        colours, and so on. Currently outputs as javascript, should also
        write to json on it's own mertits"""

        grid_width = 4

        def __init__(self, slide):
                self.utterances = [
                    ["link" for x in range(self.grid_width)]
                    for x in range(self.grid_width)]
                self.links = [
                    ["blank" for x in range(self.grid_width)]
                    for x in range(self.grid_width)]
                self.colors = [
                    ["" for x in range(self.grid_width)]
                    for x in range(self.grid_width)]
                self.tag = "unknown"
                for shape in slide.shapes:
                        self.process_shape(shape)

        def process_shape(self, shape):
                try:
                        if shape.is_placeholder:
                                if shape.placeholder_format.idx == 0:
                                        self.tag = shape.text
                        (co, ro) = Locator.get_cr(shape.top, shape.left)
                        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                                if shape.auto_shape_type == MSO_SHAPE.FOLDED_CORNER:
                                        self.links[co][ro] = "real"
                                        self.colors[co][
                                            ro] = shape.fill.fore_color.rgb
                        if shape.has_text_frame:
                                self.process_text_frame(shape, co, ro)
                except:
                        return

        def process_text_frame(self, shape, co, ro):
              #  text = self.utterances[co][ro]
                text = ""
                if "Yes" in self.utterances[co][ro]:
                        return
                for paragraph in shape.text_frame.paragraphs:
                        text += "".join([run.text.encode('ascii', 'ignore')
                                        for run in paragraph.runs])
                if text != "":
                        # add the if shape_type is text box
                        if self.links[co][ro] == "real":
                                self.links[co][ro] = make_title(text.strip())
                                self.utterances[co][ro] = "link"
                        else:
                                self.utterances[co][ro] = text.strip()

        def __str__(self):
                body = "\n".join(
                    [(self.string_from_cell(row, col))
                        for col in range(self.grid_width)
                        for row in range(self.grid_width)])
                return """
function %s(){
reset();
%s
document.main.src="ck15/CK15+.%03d.png";

}""" % (make_title(self.tag), body, slide_number)

        def string_from_cell(self, row, col):
                return '     links[{}][{}]="{}";'.format(
                    row, col, make_title(
                        self.links[row][col])) +\
                    '  utterances[{}][{}]="{}";'.format(
                        row, col, self.utterances[row][col])


def export_images(grid, slide):
        """     Second pass through shapes list finds images and saves them.
        We have to do this separately so it's guaranteed we already know what to
        name the images!"""
        images = {}
        utterances = grid.utterances
        for shape in slide.shapes:
                try:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                (co, ro) = Locator.get_cr(shape.top, shape.left)
                                if (co, ro) not in images:
                                        images[co, ro] = []
                                images[co, ro].append(shape)
                except:
                        continue

        # Compose each icon out of all the images in the grid cell.
        for (x, y) in images:
                # Go through all the images, compute bounding
                # box.
                l = min([shape.left for shape in images[x, y]])
                t = min([shape.top for shape in images[x, y]])
                r = max([shape.left +
                         shape.width for shape in images[x, y]])
                b = max([shape.top +
                         shape.height for shape in images[x, y]])

                # Scale gives us the mapping from image pixels to powerpoint
                # distance units. This depends on the resolution
                # of the images.
                scale = min([shape.width/shape.image.size[0]
                             for shape in images[x, y]])

                # Size of combined image, in actual pixels (not PPTX units)
                # If scales differ between objects, we resize
                # them next
                w = (r-l)/scale
                h = (b-t)/scale
                composite = Image.new('RGBA', (w, h))

                # Add all the images together.
                for shape in images[x, y]:
                        # TODO: flipping.
                        part = Image.open(
                            io.BytesIO(
                                shape.image.blob))
                        width = part.size[0]
                        height = part.size[1]
                        left = shape.crop_left*width
                        right = (1-shape.crop_right)*width
                        top = shape.crop_top*height
                        bottom = (1-shape.crop_bottom)*height
                        box = (int(left),
                               int(top),
                               int(right),
                               int(bottom))
                        part = part.crop(box)
                        partScale = (shape.width / part.size[0])
                        # part.size because it might have been cropped

                        part = resizeImage(part, partScale / scale)
                        composite.paste(
                            part,
                            ((shape.left - l)/scale,
                             (shape.top - t)/scale),
                            part)  # This masks out transparent pixels

                # Crop final image.
                bbox = composite.getbbox()
                composite = composite.crop(bbox)

                # Save!
                name = remove_punctuation(
                    "%d-%d-" %
                    (x, y)+utterances[x][y]) + ".png"
                folder = "icons/" + str(slide_number)
                if not os.path.exists(folder):
                        os.makedirs(folder)
                composite.save(folder + "/" + name)

prs = Presentation("../azulejoe/testSuite/launch/CK20process.pptx")
slide_number = 1
for_json = {}
for slide in prs.slides:
        grid = Grid(slide)
        for_json[slide_number] = [
            make_title(
                grid.tag),
            grid.utterances,
            grid.links,
            slide_number]
#        export_images(grid, slide)
        print grid
        slide_number += 1
#            break
with open('data.json', 'w') as outfile:
        json.dump(for_json, outfile, sort_keys=True)
# vim: tabstop=8 expandtab shiftwidth=4 softtabstop=4
